import pptx
import pdfplumber
import os

pptx_path = r'd:\BaiduSyncdisk\N8N\pauhex\AI智能分药与老年健康中枢20260113.pptx'
pdf_path = r'd:\BaiduSyncdisk\N8N\pauhex\Ai智能分药与老年健康中枢 - Prd V0内部文档20260113.pdf'
output_path = r'd:\BaiduSyncdisk\N8N\pauhex\extracted_content.md'

content = "# Extracted Content for Pauhex\n\n"

# Extract PPTX
if os.path.exists(pptx_path):
    content += "## PPTX Content\n\n"
    prs = pptx.Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        content += f"### Slide {i+1}\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
        content += "\n"

# Extract PDF
if os.path.exists(pdf_path):
    content += "## PDF Content\n\n"
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            content += f"### Page {i+1}\n"
            text = page.extract_text()
            if text:
                content += text + "\n"
            content += "\n"

with open(output_path, 'w', encoding='utf-8') as f:
    f.write(content)

print(f"Content extracted to {output_path}")
