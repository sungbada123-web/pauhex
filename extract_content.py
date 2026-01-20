
import os
import sys

def extract_from_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PPTX: {e}"

def extract_from_pdf(path):
    try:
        import PyPDF2
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PDF: {e}"

if __name__ == "__main__":
    ppt_path = r'd:\BaiduSyncdisk\N8N\pauhex\AI智能分药与老年健康中枢20260113.pptx'
    pdf_path = r'd:\BaiduSyncdisk\N8N\pauhex\Ai智能分药与老年健康中枢 - Prd V0内部文档20260113.pdf'
    
    print("--- PPTX CONTENT ---")
    print(extract_from_pptx(ppt_path))
    print("\n--- PDF CONTENT ---")
    print(extract_from_pdf(pdf_path))
